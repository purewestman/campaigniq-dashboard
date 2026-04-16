from pptx import Presentation
prs = Presentation('/Users/jwestman/Downloads/AXIZ_Path_to_100_Percent.pptx')
for i, slide in enumerate(prs.slides):
    if i > 1: break
    for shape in slide.shapes:
        if type(shape).__name__ == "Picture":
            with open(f"scratch/slide_{i}.png", "wb") as f:
                f.write(shape.image.blob)
            print(f"Saved slide_{i}.png")
