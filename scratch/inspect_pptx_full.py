from pptx import Presentation
prs = Presentation('/Users/jwestman/Downloads/AXIZ_Path_to_100_Percent.pptx')
print([len(slide.shapes) for slide in prs.slides])
