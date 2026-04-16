import sys
try:
    from pptx import Presentation
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

prs = Presentation('/Users/jwestman/Downloads/AXIZ_Path_to_100_Percent.pptx')
print(f"Total slides: {len(prs.slides)}")

for i, slide in enumerate(prs.slides):
    if i > 2: break
    print(f"\n--- Slide {i+1} ---")
    for j, shape in enumerate(slide.shapes):
        shape_type = type(shape).__name__
        text = ""
        if hasattr(shape, "text") and shape.text:
            text = shape.text[:50].replace("\n", " ")
        print(f"Shape {j}: {shape_type} | text: {text}")
        if shape_type == "Picture":
            print(f"  IMAGE! x={shape.left}, y={shape.top}, w={shape.width}, h={shape.height}")

